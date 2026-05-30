# Step 1: pip install python-pptx
# Step 2: python create_presentation.py

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

NAVY      = RGBColor(0x0D, 0x1F, 0x3C)
ORANGE    = RGBColor(0xC8, 0x74, 0x2A)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT     = RGBColor(0xAA, 0xC4, 0xE8)
DARK_CARD = RGBColor(0x1A, 0x3A, 0x6E)
MID_CARD  = RGBColor(0x0D, 0x2A, 0x55)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def set_bg(slide, color=NAVY):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def box(slide, text, l, t, w, h, sz=14, bold=False, col=WHITE, align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(sz)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = col

def rect(slide, l, t, w, h, col):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = col
    s.line.fill.background()

def header(slide, title):
    rect(slide, 0, 0, 13.33, 1.15, DARK_CARD)
    box(slide, title, 0.35, 0.15, 12.6, 0.9, sz=26, bold=True, col=ORANGE)

def bullet_slide(title, bullets, note=None):
    sl = blank(prs)
    set_bg(sl)
    header(sl, title)
    y = 1.25
    for b in bullets:
        if b.startswith("##"):
            box(sl, b[2:].strip(), 0.35, y, 12.6, 0.38, sz=13, bold=True, col=ORANGE)
            y += 0.40
        else:
            box(sl, "  \u2022  " + b, 0.35, y, 12.6, 0.40, sz=12, col=WHITE)
            y += 0.40
    if note:
        rect(sl, 0.35, 6.85, 12.6, 0.4, DARK_CARD)
        box(sl, note, 0.4, 6.87, 12.5, 0.35, sz=11, col=LIGHT, italic=True)

def table_slide(title, headers, rows, note=None):
    sl = blank(prs)
    set_bg(sl)
    header(sl, title)
    if note:
        box(sl, note, 0.35, 1.2, 12.6, 0.35, sz=11, col=LIGHT, italic=True)
    top = 1.65 if note else 1.35
    cw = 12.6 / len(headers)
    for i, h in enumerate(headers):
        rect(sl, 0.35 + i*cw, top, cw-0.04, 0.42, ORANGE)
        box(sl, h, 0.38 + i*cw, top+0.04, cw-0.08, 0.36, sz=12, bold=True, col=WHITE, align=PP_ALIGN.CENTER)
    for ri, row in enumerate(rows):
        rc = DARK_CARD if ri % 2 == 0 else MID_CARD
        for i, cell in enumerate(row):
            rect(sl, 0.35 + i*cw, top+0.42+ri*0.44, cw-0.04, 0.42, rc)
            box(sl, cell, 0.38 + i*cw, top+0.44+ri*0.44, cw-0.08, 0.38, sz=11, col=WHITE, align=PP_ALIGN.CENTER)

# SLIDE 1 — Title
s1 = blank(prs); set_bg(s1)
rect(s1, 0, 2.4, 13.33, 0.07, ORANGE)
box(s1, "PUCSL Solar & Weather Prediction System", 0.5, 0.7, 12.3, 1.1, sz=34, bold=True, col=WHITE, align=PP_ALIGN.CENTER)
box(s1, "Machine Learning-Based Renewable Energy Forecasting for Sri Lanka", 0.5, 1.9, 12.3, 0.55, sz=18, col=LIGHT, align=PP_ALIGN.CENTER, italic=True)
box(s1, "Research Project 2026  |  Public Utilities Commission of Sri Lanka", 0.5, 2.9, 12.3, 0.45, sz=14, bold=True, col=ORANGE, align=PP_ALIGN.CENTER)
box(s1, "15 Years Data (2010-2025)  |  1,242 Grid Cells  |  475 Real Solar Plants  |  4 Prediction Equations", 0.5, 5.9, 12.3, 0.45, sz=12, col=LIGHT, align=PP_ALIGN.CENTER)

# SLIDE 2 — Overview
bullet_slide("Project Overview", [
    "Sri Lanka needs accurate renewable energy forecasting for national grid stability",
    "This system predicts Solar Generation, Temperature & Rainfall for any location",
    "Entire country covered using a 10km x 10km grid = 1,242 grid cells",
    "Built using 15 years of historical weather data (2010-2025)",
    "##Data Sources:",
    "OpenMeteo API - reanalysis weather data for all 1,242 grid cells",
    "24 official Meteorological Stations (Department of Meteorology, Sri Lanka)",
    "475 real solar plants (PUCSL) - 8,201 rows of monthly generation data",
    "##Final Output:",
    "Monthly & annual solar generation (kWh) for any location and installed capacity (kWp)"])

# SLIDE 3 — Pipeline
bullet_slide("Data Pipeline - How It All Flows", [
    "Step 1:  Raw weather downloaded from OpenMeteo API -> T_om, RF_om for 1,242 cells",
    "Step 2:  Actual measurements from 24 meteorological stations collected",
    "Step 3:  Error Correction Equations applied -> T_corrected, RF_corrected",
    "Step 4:  Corrected data merged with real solar plant generation records",
    "Step 5:  Machine Learning models trained (Temperature, Rainfall, Solar)",
    "Step 6:  Models applied to all 1,242 grid cells x 12 months = 14,904 predictions",
    "Step 7:  Results visualized on interactive maps and Streamlit dashboard"])

# SLIDE 4 — Beta
bullet_slide("What are Beta (B) Coefficients?", [
    "Beta values are the coefficients the model LEARNS from historical training data",
    "They represent how much each input variable contributes to the prediction",
    "##Beta symbols and meaning:",
    "B0  =  Intercept - base value when all inputs equal zero",
    "B1, B2, B3 ...  =  Weights assigned to each input variable",
    "##Interpreting Beta values:",
    "Large positive B  ->  that variable strongly INCREASES the output",
    "Large negative B  ->  that variable strongly DECREASES the output",
    "B close to 0  ->  that variable has LITTLE effect on the output",
    "##How Beta is calculated - Ordinary Least Squares (OLS):",
    "Minimize:  Sum of ( yi - y_predicted )^2   where yi = actual,  y_predicted = predicted"])

# SLIDE 5 — sin/cos encoding
bullet_slide("Why sin(m) and cos(m)? - Cyclical Month Encoding", [
    "Raw month number (1-12) treats December(12) and January(1) as FAR apart numerically",
    "But December and January are climatically ADJACENT - this causes prediction errors!",
    "##Solution: Encode month as sine and cosine to wrap the calendar into a smooth circle:",
    "month_sin  =  sin( 2*pi * month / 12 )",
    "month_cos  =  cos( 2*pi * month / 12 )",
    "##Example values:",
    "January   ->  sin = +0.500,  cos = +0.866",
    "July      ->  sin = -0.500,  cos = -0.866",
    "December  ->  sin = -0.866,  cos = +0.500  (close to January - correct!)",
    "##Interaction features also used:",
    "lat x sin(m),  lat x cos(m),  lon x sin(m),  lon x cos(m),  elev x sin(m),  elev x cos(m)"])

# SLIDE 6 — Equation 1
bullet_slide("Equation 1 - Temperature Error Correction", [
    "##Purpose:",
    "OpenMeteo reanalysis data has systematic bias vs actual ground measurements",
    "This equation corrects that bias using official met station data",
    "##Equation:",
    "T_corrected = B0 + B1.lat + B2.lon + B3.elev + B4.month + B5.T_om",
    "##Input Variables:",
    "lat = Latitude (degrees N)    |    lon = Longitude (degrees E)",
    "elev = Elevation in meters above sea level",
    "month = Month number (1 = January ... 12 = December)",
    "T_om = Raw temperature from OpenMeteo (degrees C)",
    "##Output:  T_corrected = Bias-corrected temperature (degrees C)",
    "B3 (elevation) is NEGATIVE - higher elevation = cooler temperature",
    "B5 (T_om) close to 1.0 - OpenMeteo is mostly accurate, small correction needed"])

# SLIDE 7 — Equation 2
bullet_slide("Equation 2 - Rainfall Error Correction", [
    "##Purpose:  Correct OpenMeteo rainfall bias using actual met station measurements",
    "##Equation:",
    "ln(RF_corrected) = B0 + B1.lat + B2.lon + B3.elev + B4.month + B5.ln(RF_om + 1)",
    "##Input Variables:",
    "lat, lon, elev, month = same as temperature equation",
    "RF_om = Raw rainfall from OpenMeteo (mm)",
    "ln(RF_om + 1) = Natural log of raw rainfall  (+1 added to handle zero-rainfall months)",
    "##Output:  RF_corrected = Bias-corrected rainfall (mm)",
    "##Why log transformation?",
    "Rainfall is right-skewed - log normalizes the distribution for better accuracy",
    "+1 inside the log prevents ln(0) which is mathematically undefined"])

# SLIDE 8 — Temp Model Features
bullet_slide("Temperature Prediction Model - Features & Data", [
    "##Purpose:  Predict monthly temperature for any location and year in Sri Lanka",
    "##12 Input Features:",
    "lat, lon, elevation, month_sin, month_cos, year",
    "lat x sin(m), lat x cos(m), lon x sin(m), lon x cos(m), elev x sin(m), elev x cos(m)",
    "##Training Data:",
    "44,047 rows - 207 OpenMeteo grid cells + 24 meteorological stations merged",
    "Latitude range: 5.96 N to 9.68 N (full Sri Lanka including northern cities)",
    "Years covered: 2010-2025 (15 years)"])

# SLIDE 9 — Temp Model Results
table_slide("Temperature Model - Results Comparison",
    ["Model", "R2 Score", "MAE", "Verdict"],
    [["Linear Regression",     "0.8528", "0.6799 C", "Good"],
     ["Polynomial (degree 2)", "0.9054", "0.5341 C", "Better"],
     ["Random Forest  BEST",   "0.9870", "0.1769 C", "Excellent"],
     ["Gradient Boosting",     "0.9619", "0.3403 C", "Very Good"]],
    note="R2=0.987 means 98.7% of temperature variance explained  |  MAE = only 0.18 C average error")

# SLIDE 10 — Rainfall Model
bullet_slide("Rainfall Prediction Model", [
    "##Purpose:  Predict monthly rainfall for any location and year in Sri Lanka",
    "##Same 12 input features as temperature model:",
    "lat, lon, elevation, month_sin, month_cos, year + 6 interaction features",
    "##Best Model:  Random Forest  -  R2 = 0.956",
    "95.6% of rainfall variance is explained by the model",
    "##Why slightly lower R2 than temperature?",
    "Rainfall is more spatially variable and event-driven than temperature",
    "A single storm can cause large local variation hard to capture at grid level",
    "Nevertheless 95.6% is excellent accuracy for rainfall prediction",
    "##Training Data:",
    "Same 44,047 rows covering full Sri Lanka (5.96 N to 9.68 N)"])

# SLIDE 11 — Equation 3
bullet_slide("Equation 3 - Solar Specific Yield Prediction", [
    "##Purpose:  Predict solar energy per kWp of installed capacity per month",
    "##Equation:",
    "SpecificYield = B0 + B1.lat + B2.lon + B3.elev + B4.sin(m) + B5.cos(m) + B6.T_corrected + B7.RF_corrected",
    "##Coefficient Interpretation:",
    "B1 (lat): NEGATIVE - higher latitude = less solar yield",
    "B2 (lon): NEGATIVE - moving east = wetter climate = less solar",
    "B3 (elev): POSITIVE - higher elevation = clearer atmosphere = more solar",
    "B6 (T): POSITIVE - warmer = more sunshine hours = more solar",
    "B7 (RF): NEGATIVE - more rain = cloud cover = less solar",
    "##Fitted equation from actual training data:",
    "SpecificYield = 769.26 - 4.20.lat - 8.03.lon + 0.008.elev - 2.07.sin(m) + 1.17.cos(m) + 1.59.T - 8.72.ln(RF)"])

# SLIDE 12 — Equation 4
bullet_slide("Equation 4 - Final Energy Generation Prediction", [
    "Final step: converting specific yield into actual energy generation for a plant",
    "##Equation:",
    "Predicted kWh  =  Specific Yield (kWh/kWp)  x  Capacity (kWp)",
    "##Worked Example:",
    "Location: Colombo  |  Month: February  |  Installed Capacity: 10 kWp",
    "Specific Yield predicted by model = 110 kWh/kWp",
    "Predicted Generation  =  110  x  10  =  1,100 kWh for February",
    "Annual Generation  =  Sum of all 12 monthly predictions",
    "##Why this matters for PUCSL:",
    "Any plant size can be evaluated - from 5 kWp school rooftop to 1,000 kWp solar farm",
    "Grid planners can estimate contribution before installation begins"])

# SLIDE 13 — Solar Model Results
table_slide("Solar Generation Model - Results Comparison",
    ["Model", "R2 Score", "MAE", "Notes"],
    [["Linear Regression (log RF)",    "0.0799", "25.28 kWh/kWp", "Baseline"],
     ["Linear Regression (plain RF)",  "0.1160", "24.46 kWh/kWp", "Better"],
     ["Polynomial Regression deg 2",   "0.1240", "24.38 kWh/kWp", "Best - used"]],
    note="Lower R2 expected: solar also depends on panel brand, inverter efficiency, shading, dust - not in dataset")

# SLIDE 14 — Grid Results
table_slide("Sri Lanka Solar Potential Mapping - Results",
    ["Category", "Annual Yield (kWh/kWp)", "Grid Cells"],
    [["Very High", "> 1,320",        "1"],
     ["High",      "1,200 - 1,320",  "226"],
     ["Medium",    "1,000 - 1,200",  "484"],
     ["Low",       "< 1,000",        "345"]],
    note="1,242 cells x 12 months = 14,904 predictions  |  Top: Colombo, Moratuwa  |  Lowest: Near Mullaitivu")

# SLIDE 15 — Correlations
table_slide("Correlation Analysis - Key Findings",
    ["Relationship", "r Value", "Strength", "Physical Meaning"],
    [["Solar vs Temperature",    "r = +0.50", "Moderate Positive",    "Warmer -> more sun -> more solar"],
     ["Solar vs Rainfall",       "r = -0.71", "Strong Negative",      "More rain -> clouds -> less solar"],
     ["Temperature vs Rainfall", "r = -0.91", "Very Strong Negative", "Hot dry vs cool wet months"]],
    note="These correlations confirm that model features are physically meaningful and correctly chosen")

# SLIDE 16 — Dashboard
bullet_slide("Streamlit Dashboard - System Features", [
    "A fully interactive web application built with Python Streamlit",
    "##Dashboard Modules:",
    "Home - Project overview, key statistics, institutional summary",
    "Solar Prediction - Monthly/annual kWh for any city, capacity and year",
    "Weather Prediction - Monthly temperature and rainfall for any location",
    "Maps Viewer - Interactive Folium maps: solar, temperature, rainfall across Sri Lanka",
    "Correlation Analysis - r values between solar, temperature, and rainfall",
    "Data Visualizations - Spatial distribution charts by district and province",
    "##Technical Stack:",
    "Python  |  Streamlit  |  Scikit-learn  |  Pandas  |  Folium  |  Plotly  |  Joblib"])

# SLIDE 17 — Conclusion
bullet_slide("Key Findings and Contributions", [
    "Complete ML pipeline built for renewable energy prediction across Sri Lanka",
    "##Model Performance:",
    "Temperature model:  R2 = 0.987  - near-perfect accuracy (MAE = 0.18 C)",
    "Rainfall model:     R2 = 0.956  - excellent accuracy",
    "Solar model:  Polynomial regression - physically interpretable equation for PUCSL",
    "##Coverage:",
    "All 1,242 grid cells across Sri Lanka mapped with solar potential",
    "Northern Sri Lanka (Jaffna, Vavuniya, Mannar) correctly predicted with station data",
    "15 years of historical data (2010-2025) ensures robust and reliable predictions",
    "##Impact:",
    "PUCSL can now forecast solar contributions before any plant is installed",
    "Supports national energy transition planning and grid stability management"])

# Save
os.makedirs("output", exist_ok=True)
prs.save("output/PUCSL_Solar_Weather_Prediction.pptx")
print(f"Done! Saved to output/PUCSL_Solar_Weather_Prediction.pptx")
print(f"Total slides: {len(prs.slides)}")
